from openai import OpenAI
import tiktoken
import secret_vars
import json
import prompts
import time
import random
import string
import os, shutil
import hashlib
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from docx import Document
from urllib.parse import unquote


FOLDER_PATH = os.getcwd()
if os.getcwd() == "/":
    FOLDER_PATH = "/home/maxet24/doklad_gen"

WEB_HEADERS = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/103.0.5060.114 Safari/537.36"
    }

MAX_SINGLE_IMG_WEIGHT = 300000
DEL_IMGS = True
DEBUG = False
GPT_MODEL = "gpt-3.5-turbo"
client = OpenAI(
    api_key=secret_vars.API_KEY,
)

start_time = time.time()


def get_completion(messages, model=GPT_MODEL):
    response = client.chat.completions.create(
        model=model,
        messages=messages,
        temperature=0
    )

    cost = 0
    if GPT_MODEL == "gpt-3.5-turbo-16k":
        cost = response.usage.total_tokens / 1000 * 0.003
    if GPT_MODEL == "gpt-3.5-turbo":
        cost = response.usage.total_tokens / 1000 * 0.001

    if DEBUG:
        print(cost, '$ COST')

    return response.choices[0].message.content, cost

def get_report_plan(theme, prompt_type = "new"):
    plan_prompt = f"""Составь план доклада для университета на тему: {theme}"""
    if prompt_type == "old":
        plan_prompt += prompts.plan_example_old
    else:
        plan_prompt += prompts.plan_example
    # plan_prompt = f"""Напиши фрагмент доклада:"""

    messages = [
        {"role": "user", "content": plan_prompt}
    ]

    resp, money_spent = get_completion(messages)
    resp = json.loads(resp)

    return resp, money_spent

def timer_start():
    start_time = time.time()

def timer_check():
    return time.time() - start_time

def generate_random_hash():
    """Generate a random hash using sha256."""
    characters = string.ascii_letters + string.digits
    random_string = ''.join(random.choice(characters) for i in range(10))
    hash_object = hashlib.sha256(random_string.encode())
    return hash_object.hexdigest()

def gen_doklad_pptx(theme, fio, filename):
    timer_start()
    total_cost = 0

    doklad = {}

    parts, money_spent_plan = get_report_plan(theme, "old")
    total_cost += money_spent_plan
    # parts = {'Введение': ['Общая информация о ледниках', 'Значение изучения происхождения и распространения ледников'], 'Раздел 1: Происхождение ледников': ['Геологические факторы, влияющие на формирование ледников', 'Роль климатических изменений в происхождении ледников', 'Теории происхождения ледников'], 'Раздел 2: Распространение ледников на земле': ['Географическое распределение ледников', 'Факторы, влияющие на распространение ледников', 'Изменения в распространении ледников в истории Земли'], 'Раздел 3: Основные типы ледников': ['Континентальные ледники', 'Альпийские ледники', 'Побережные ледники', 'Ледники на вулканах'], 'Раздел 4: Строение ледников': ['Зоны ледника: верхняя, средняя и нижняя', 'Ледовые языки и ледовые шапки', 'Ледовые трещины и сераки', 'Движение ледников'], 'Заключение': ['Выводы о происхождении и распространении ледников', 'Значение изучения ледников для науки и практического применения']}

    for part in parts.keys():
        doklad[part] = []
        if DEBUG:
            print('------------')
            print(part)
        for pod_part in parts[part]:
            if DEBUG:
                print(pod_part)
                print("Генерирую...")
            punkt_prompt = f"""Тебе нужно написать часть доклада (не весь) на тему \"{theme}\"\
в пункте \"{pod_part}\" из раздела  \"{part}\"
Пиши только самую важную информацию, которая относится к пункту: \"{pod_part}\"
Объем: 50 слов"""
            messages = [
                {"role": "user", "content": punkt_prompt}
            ]
            resp_raw = get_completion(messages)
            resp = resp_raw[0]
            total_cost += resp_raw[1]

            if DEBUG:
                print(resp)
            doklad[part].append({pod_part:resp})

    if DEBUG:
        print(f"{timer_check()} seconds.")
        print(f'{total_cost} $')

    create_presentation_from_json(doklad, theme, fio, filename)

    if DEBUG:
        print(f"{timer_check()} seconds.")

    meta = {
        'money_spent': total_cost
    }
    return meta

def gen_doklad_pptx_docx(theme, fio, pptx_filename, docs_filename):
    timer_start()
    total_cost = 0

    doklad_dict = {}

    parts, money_spent_plan = get_report_plan(theme)


    # parts = {'Введение': ['Общая информация о ледниках', 'Значение изучения происхождения и распространения ледников'], 'Раздел 1: Происхождение ледников': ['Геологические факторы, влияющие на формирование ледников', 'Роль климатических изменений в происхождении ледников', 'Теории происхождения ледников'], 'Раздел 2: Распространение ледников на земле': ['Географическое распределение ледников', 'Факторы, влияющие на распространение ледников', 'Изменения в распространении ледников в истории Земли'], 'Раздел 3: Основные типы ледников': ['Континентальные ледники', 'Альпийские ледники', 'Побережные ледники', 'Ледники на вулканах'], 'Раздел 4: Строение ледников': ['Зоны ледника: верхняя, средняя и нижняя', 'Ледовые языки и ледовые шапки', 'Ледовые трещины и сераки', 'Движение ледников'], 'Заключение': ['Выводы о происхождении и распространении ледников', 'Значение изучения ледников для науки и практического применения']}

    for part in parts.keys():
        doklad_dict[part] = []
        if DEBUG:
            print('-----МЕЖДУ РАЗДЕЛАМИ------')
            print(part)
        for pod_part in parts[part]:
            if DEBUG:
                print(f"Генерирую раздел '{pod_part}'...")
            punkt_prompt = f"""Тебе нужно написать часть доклада (не весь) на тему \"{theme}\"\
 в пункте \"{pod_part[0]}\" из раздела  \"{part}\"
 Пиши только самую важную информацию, которая относится к пункту: \"{pod_part[0]}\"
 Объем: 50 слов"""
            messages = [
                {"role": "user", "content": punkt_prompt}
            ]

            resp_raw = get_completion(messages)
            resp = resp_raw[0]
            total_cost += resp_raw[1]

            if DEBUG:
                print(resp)
            doklad_dict[part].append({
                pod_part[0]: {
                    'text': resp,
                    'img_query': pod_part[1]
                }
            })

    if DEBUG:
        print(f"{timer_check()} seconds.")
        print(f'{total_cost} $')
    
    print(doklad_dict)

    create_word_from_json(doklad_dict, theme, fio, docs_filename)
    create_presentation_with_images_from_json(doklad_dict, theme, fio, pptx_filename)

    meta = {
        'money_spent': total_cost
    }
    return meta

def add_background_image(presentation, section_slide, background_image_path):
    background_picture = section_slide.shapes.add_picture(background_image_path, 0, 0, presentation.slide_width,
                                                          presentation.slide_height)
    section_slide.shapes._spTree.remove(background_picture._element)
    section_slide.shapes._spTree.insert(2, background_picture._element)

def add_foreground_image(presentation, section_slide, background_image_path):
    IMAGE_WIDTH = presentation.slide_width / 1.5
    foreground_picture = section_slide.shapes.add_picture(background_image_path, (presentation.slide_width - IMAGE_WIDTH) / 2, presentation.slide_height / 3, IMAGE_WIDTH)
    section_slide.shapes._spTree.remove(foreground_picture._element)
    section_slide.shapes._spTree.insert(2, foreground_picture._element)

def add_word_section(doc, title, content):
    doc.add_heading(title, level=1)
    for item in content:
        for subtitle, punkt_content in item.items():
            doc.add_heading(subtitle, level=2)
            doc.add_paragraph(punkt_content['text'])

def create_word_from_json(json_data, theme, fio, filename):
    doc = Document()

    for section_title, section_content in json_data.items():
        add_word_section(doc, section_title, section_content)

    doc.save(filename)

def save_html_to_file(html_content, file_path):
    try:
        with open(file_path, 'w', encoding='utf-8') as file:
            file.write(html_content)
        print(f"HTML content saved to {file_path}")
    except Exception as e:
        print(f"An error occurred while saving the HTML content: {e}")

def get_html_by_query(query):

    params = {
        "q": query,  # search query
        "tbm": "isch",  # image results
        "hl": "ru",  # language of the search
        "gl": "us",  # country where search comes from
        "ijn": "0"  # page number
    }

    try:
        response = requests.get("https://www.google.com/search", params=params, headers=WEB_HEADERS, timeout=30)
        # response = requests.get(url)
        response.raise_for_status()

        return response.text
    except requests.exceptions.RequestException as e:
        return '-1'

def is_image_url(url):
    try:
        response = requests.head(url)
        content_type = response.headers.get('content-type')
        # is accessable and weights normal (< 300 kb)
        if DEBUG:
            print(response.headers.get('Content-Length') + " bytes")
        return content_type.startswith('image') and int(response.headers.get('Content-Length')) < MAX_SINGLE_IMG_WEIGHT
    except:
        return False

def get_image_url_by_query(query):
    # КРИВО, КОСО, НО РАБИТ
    html_text = get_html_by_query(query)
    url = ''

    while 'http' not in url or not is_image_url(url):
        first_jpeg = html_text.index(".jpg")

        first_quote = -1
        second_quote = -1
        for i in range(first_jpeg, len(html_text)):
            if html_text[i] == '\"':
                second_quote = i
                break

        for i in range(first_jpeg, 0, -1):
            if html_text[i] == '\"':
                first_quote = i
                break

        url = html_text[first_quote + 1: second_quote]
        html_text = html_text.replace(html_text[first_quote + 1: second_quote], '')

    return url

def get_image_by_query(query, file_path):
    url = get_image_url_by_query(query)
    if DEBUG:
        print(f'Got an URL: {url}')

    is_success = False
    for i in range(10):
        try:
            response = requests.get(url, headers=WEB_HEADERS)
            response.raise_for_status()

            with open(file_path, 'wb') as file:
                file.write(response.content)
            if DEBUG:
                print(f"Image downloaded and saved to {file_path}")

            is_success = True

        except requests.exceptions.RequestException as e:
            print(f"An error occurred: {e}\n'{query}': {url}")

        if is_success:
            break
        else:
            print('Fail, another try, i = ' + str(i))

def create_presentation_from_json(json_data, theme, fio, filename):

    presentation = Presentation()

    titul_slide_layout = presentation.slide_layouts[6]
    titul_slide = presentation.slides.add_slide(titul_slide_layout)

    # Title
    width = Inches(8)
    height = Inches(2)
    section_title = titul_slide.shapes.add_textbox((presentation.slide_width - width) / 2,
                                                   (presentation.slide_height - height) / 2, width, height)
    text_frame = section_title.text_frame
    text_frame.word_wrap = True
    p_title = text_frame.paragraphs[0]
    p_title.text = theme
    p_title.alignment = PP_ALIGN.CENTER
    p_title.font.size = Pt(32)

    # FIO
    section_fio = titul_slide.shapes.add_textbox(left=Inches(6.5), top=Inches(6.8), width=Inches(3), height=Inches(0.5))
    fio_frame = section_fio.text_frame
    p_fio = fio_frame.paragraphs[0]
    p_fio.text = fio
    p_fio.font.size = Pt(16)
    p_fio.alignment = PP_ALIGN.RIGHT

    add_background_image(presentation, titul_slide, FOLDER_PATH + "/images/listya.jpg")

    for section, slides_data in json_data.items():

        # Добавляем слайд с заголовком раздела
        section_slide_layout = presentation.slide_layouts[6]
        section_slide = presentation.slides.add_slide(section_slide_layout)

        # Title
        width = Inches(8)
        height = Inches(2)
        section_title = section_slide.shapes.add_textbox((presentation.slide_width - width) / 2,
                                                       (presentation.slide_height - height) / 2, width, height)
        text_frame = section_title.text_frame
        text_frame.word_wrap = True
        p_title = text_frame.paragraphs[0]
        p_title.text = section
        p_title.alignment = PP_ALIGN.CENTER
        p_title.font.size = Pt(32)

        add_background_image(presentation, section_slide, FOLDER_PATH + "/images/listya.jpg")

        # Добавляем слайды с пунктами в разделе
        for slide_data in slides_data:

            bullet_slide_layout = presentation.slide_layouts[1]
            bullet_slide = presentation.slides.add_slide(bullet_slide_layout)

            add_background_image(presentation, bullet_slide, FOLDER_PATH + "/images/luzhy.jpg")

            title = bullet_slide.shapes.title
            content = bullet_slide.placeholders[1]


            for slide_title, slide_content in slide_data.items():
                title.text = slide_title
                content.text = slide_content
            for parag in content.text_frame.paragraphs:
                parag.font.size = Pt(26)
            title.text_frame.paragraphs[0].font.size = Pt(40)


    presentation.save(filename)

def create_presentation_with_images_from_json(json_data, theme, fio, filename):
    # Generate hash for session
    request_hash = generate_random_hash()
    # create new folder
    os.mkdir(FOLDER_PATH + "/doklads/" + request_hash)

    print()
    # Get images
    img_paths = {}
    img_num = 1
    for section, slides_data in json_data.items():
        for slide_data in slides_data:
            for slide_title, slide_content in slide_data.items():
                img_path = FOLDER_PATH + "/doklads/" + request_hash + '/' + str(img_num) + ".jpeg"
                img_paths[slide_content['img_query']] = img_path
                if DEBUG:
                    print(img_path, slide_content['img_query'])
                get_image_by_query(slide_content['img_query'], img_path)
                img_num += 1

    # Create presentation
    presentation = Presentation()

    titul_slide_layout = presentation.slide_layouts[6]
    titul_slide = presentation.slides.add_slide(titul_slide_layout)

    # Title
    width = Inches(8)
    height = Inches(2)
    section_title = titul_slide.shapes.add_textbox((presentation.slide_width - width) / 2,
                                                   (presentation.slide_height - height) / 2, width, height)
    text_frame = section_title.text_frame
    text_frame.word_wrap = True
    p_title = text_frame.paragraphs[0]
    p_title.text = theme
    p_title.alignment = PP_ALIGN.CENTER
    p_title.font.size = Pt(32)

    # FIO
    section_fio = titul_slide.shapes.add_textbox(left=Inches(6.5), top=Inches(6.8), width=Inches(3), height=Inches(0.5))
    fio_frame = section_fio.text_frame
    p_fio = fio_frame.paragraphs[0]
    p_fio.text = fio
    p_fio.font.size = Pt(16)
    p_fio.alignment = PP_ALIGN.RIGHT

    add_background_image(presentation, titul_slide, FOLDER_PATH + "/images/listya.jpg")

    for section, slides_data in json_data.items():

        # Добавляем слайд с заголовком раздела
        section_slide_layout = presentation.slide_layouts[6]
        section_slide = presentation.slides.add_slide(section_slide_layout)

        # Title
        width = Inches(8)
        height = Inches(2)
        section_title = section_slide.shapes.add_textbox((presentation.slide_width - width) / 2,
                                                         (presentation.slide_height - height) / 2, width, height)
        text_frame = section_title.text_frame
        text_frame.word_wrap = True
        p_title = text_frame.paragraphs[0]
        p_title.text = section
        p_title.alignment = PP_ALIGN.CENTER
        p_title.font.size = Pt(32)

        add_background_image(presentation, section_slide, FOLDER_PATH + "/images/listya.jpg")

        # Добавляем слайды с пунктами в разделе
        for slide_data in slides_data:

            bullet_slide_layout = presentation.slide_layouts[6]
            bullet_slide = presentation.slides.add_slide(bullet_slide_layout)

            # Title
            width = Inches(8)
            height = Inches(2)
            section_title = bullet_slide.shapes.add_textbox((presentation.slide_width - width) / 2,
                                                             Inches(0.1), width, height)
            text_frame = section_title.text_frame
            text_frame.word_wrap = True
            p_title = text_frame.paragraphs[0]
            p_title.alignment = PP_ALIGN.CENTER
            p_title.font.size = Pt(32)

            for slide_title, slide_content in slide_data.items():
                try:
                    add_foreground_image(presentation, bullet_slide, img_paths[slide_content['img_query']])
                    add_background_image(presentation, bullet_slide, FOLDER_PATH + "/images/luzhy.jpg")
                except:
                    print(f'IMG NOT FOUND: {img_paths[slide_title]}\n{slide_title}')
                p_title.text = slide_title

    presentation.save(filename)

    # delete images
    if DEL_IMGS:
        shutil.rmtree(FOLDER_PATH + "/doklads/" + request_hash + '/')


if __name__ == "__main__":
    fio = "Глебов Максим Александрович РСБО-01-23"
    THEME = "Происхождение и распространение на земле ледников. Основные типы и строение ледников."
    if DEBUG:
        print("Process started...")
    TEST_PATH = "D:\ProgrammingProjects\doklad_gen/doklads/test.jpeg"
    # url = 'https://ya.ru/images/search?from=tabbar&text=%D0%BA%D0%BE%D1%82%D1%8B'
    # get_image_url_by_query(url)

    # russian_words = [
    #     "яблоко", "молоко", "солнце", "дом", "книга", "город", "музыка",
    #     "программирование", "вечер", "зима", "весна", "лето", "осень", "кот",
    #     "собака", "стол", "стул", "красивый", "интересный", "успех", "компьютер",
    #     "дружба", "семья", "учеба", "работа", "здоровье", "спорт", "отдых", "путешествие"
    # ]
    #
    # for i in range(len(russian_words)):
    #     get_image_by_query(russian_words[i], f'00{i}.jpg')

    # print(get_image_by_query("ледниковые исследования", TEST_PATH))

    # print(gen_doklad_pptx_docx(THEME, fio, 'test.pptx', 'test.docs'))

    # create_presentation_with_images_from_json(prompts.test_doklad_with_img, THEME, fio, "test.pptx")

    # create_word_from_json(prompts.test_doklad, THEME, fio, f"{FOLDER_PATH}/doklads/{'maxet24_1.docx'}")

    gen_doklad_pptx_docx(THEME, fio, f"{FOLDER_PATH}/doklads/{'maxet24_1.pptx'}", f"{FOLDER_PATH}/doklads/{'maxet24_1.docx'}")

    # print(is_image_url("https://www.hse.ru/data/2020/12/11/1356352199/35815118_1818082964897249_7399175503552184320_o.jpg"))

    # create_presentation_with_images_from_json(prompts.test_doklad_with_img, THEME, fio, 'test.pptx')




